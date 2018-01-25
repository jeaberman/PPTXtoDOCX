from pptx import Presentation
import os
from docx import Document
import pathlib

for f in os.listdir(os.curdir):
	if f.endswith(".pptx"):
		filepath = f
		filename = f

document= Document()
document.add_heading(filename,0)

prs = Presentation(filepath)

for slide in prs.slides:
	slidenum = 'Slide ' + str(prs.slides.index(slide)+ 1) + ':'
	document.add_heading(slidenum,1)
	for shape in slide.shapes:
		if not shape.has_text_frame:
			continue
		for paragraph in shape.text_frame.paragraphs:
			for run in paragraph.runs:
				document.add_paragraph(run.text)
	if slide.has_notes_slide:
		document.add_heading('Notes:',2)
		notes = slide.notes_slide.notes_text_frame
		document.add_paragraph(notes.text)

document.save('Output.docx')